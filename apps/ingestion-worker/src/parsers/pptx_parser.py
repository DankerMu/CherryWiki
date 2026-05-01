from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_parser import BaseParser, ParseResult
from .utils import compact_blank_lines, package_version


class PptxParser(BaseParser):
    source_type = "pptx"
    extraction_tool = "python-pptx"
    extraction_version = package_version("python-pptx")

    def parse(self, file_path: Path) -> ParseResult:
        presentation = Presentation(file_path)
        sections: list[str] = []
        image_count = 0

        for index, slide in enumerate(presentation.slides, start=1):
            title = (
                slide.shapes.title.text.strip()
                if slide.shapes.title and slide.shapes.title.has_text_frame
                else ""
            )
            sections.append(
                f"## Slide {index}: {title}" if title else f"## Slide {index}"
            )

            seen_title = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                if title and text == title and not seen_title:
                    seen_title = True
                    continue
                sections.append(text)

            notes = self._speaker_notes(slide)
            if notes:
                sections.append(
                    "\n".join(
                        f"> {line}" if line else ">" for line in notes.splitlines()
                    )
                )

        content = compact_blank_lines("\n\n".join(sections))
        return ParseResult(
            content=content,
            metadata={
                "source_type": self.source_type,
                "extraction_tool": self.extraction_tool,
                "extraction_version": self.extraction_version,
                "extraction_params": {"extract_notes": True},
                "page_count": len(presentation.slides),
                "char_count": len(content),
                "image_count": image_count,
            },
        )

    def _speaker_notes(self, slide: object) -> str:
        try:
            notes_slide = slide.notes_slide
            frame = notes_slide.notes_text_frame
        except Exception:
            return ""
        if frame is None:
            return ""
        return "\n".join(
            paragraph.text.strip()
            for paragraph in frame.paragraphs
            if paragraph.text.strip()
        )
