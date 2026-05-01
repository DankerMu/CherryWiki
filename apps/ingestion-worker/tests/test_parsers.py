from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation

from src.parsers.docx_parser import DocxParser
from src.parsers.image_parser import ImageParser
from src.parsers.pdf_parser import PdfParser
from src.parsers.pptx_parser import PptxParser
from src.parsers.text_parser import TextParser
from src.parsers.xlsx_parser import XlsxParser


def test_text_parser_reads_markdown(tmp_path: Path) -> None:
    path = tmp_path / "note.md"
    path.write_text("# Title\n\nBody text", encoding="utf-8")

    result = TextParser().parse(path)

    assert "# Title" in result.content
    assert result.metadata["extraction_tool"] == "passthrough"
    assert result.metadata["page_count"] == 1


def test_pdf_parser_extracts_text(tmp_path: Path) -> None:
    path = tmp_path / "sample.pdf"
    path.write_bytes(_minimal_pdf("Hello PDF"))

    result = PdfParser(ocr_enabled=False).parse(path)

    assert "## Page 1" in result.content
    assert "Hello PDF" in result.content
    assert result.metadata["page_count"] == 1


def test_docx_parser_extracts_headings_paragraphs_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("Doc Heading", level=1)
    document.add_paragraph("Paragraph body")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    document.save(path)

    result = DocxParser().parse(path)

    assert "# Doc Heading" in result.content
    assert "Paragraph body" in result.content
    assert "| A | B |" in result.content


def test_pptx_parser_extracts_slide_text(tmp_path: Path) -> None:
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide body"
    presentation.save(path)

    result = PptxParser().parse(path)

    assert "## Slide 1: Slide Title" in result.content
    assert "Slide body" in result.content
    assert result.metadata["page_count"] == 1


def test_xlsx_parser_extracts_sheet_tables(tmp_path: Path) -> None:
    path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Name", "Count"])
    sheet.append(["Cherry", 3])
    workbook.save(path)

    result = XlsxParser().parse(path)

    assert "## Sheet: Data" in result.content
    assert "| Name | Count |" in result.content
    assert "| Cherry | 3 |" in result.content


def test_image_parser_returns_metadata_without_ocr(tmp_path: Path) -> None:
    path = tmp_path / "image.png"
    Image.new("RGB", (20, 10), color="white").save(path)

    result = ImageParser().parse(path)

    assert "# Image OCR Result" in result.content
    assert "Dimensions: 20x10" in result.content
    assert result.metadata["image_count"] == 1


def _minimal_pdf(text: str) -> bytes:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 24 Tf 72 720 Td ({escaped}) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        f"<< /Length {len(stream.encode('latin-1'))} >>\nstream\n{stream}\nendstream".encode(
            "latin-1"
        ),
    ]
    body = b"%PDF-1.4\n"
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(body))
        body += f"{index} 0 obj\n".encode("ascii") + obj + b"\nendobj\n"
    xref_offset = len(body)
    xref = [b"xref\n0 6\n", b"0000000000 65535 f \n"]
    xref.extend(f"{offset:010d} 00000 n \n".encode("ascii") for offset in offsets[1:])
    body += b"".join(xref)
    body += b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
    body += f"startxref\n{xref_offset}\n%%EOF\n".encode("ascii")
    return body
